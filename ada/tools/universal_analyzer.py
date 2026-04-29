#!/usr/bin/env python3
"""
Universal Analyzer Tool for Ada
Extracts text from PDFs and Office files (DOC, DOCX, XLS, XLSX, PPT, PPTX).
"""

import os
import sys
import argparse
import logging
from pathlib import Path

# Try to import optional dependencies
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    if PyPDF2 is None:
        raise ImportError("PyPDF2 not installed. Run: pip install PyPDF2")
    
    text_parts = []
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
    return '\n\n'.join(text_parts)


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    if Document is None:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    
    doc = Document(file_path)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    return '\n\n'.join(text_parts)


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file."""
    if openpyxl is None:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")
    
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"=== Sheet: {sheet_name} ===")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else '' for cell in row]
            if any(row_values):
                text_parts.append('\t'.join(row_values))
    return '\n'.join(text_parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    if Presentation is None:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            text_parts.append(f"--- Slide {slide_num} ---\n" + '\n'.join(slide_text))
    return '\n\n'.join(text_parts)


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        return file.read()


def analyze_file(file_path: str) -> dict:
    """Analyze file and extract text based on extension."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    
    ext = path.suffix.lower()
    result = {
        'file_path': str(path),
        'file_name': path.name,
        'file_size': path.stat().st_size,
        'extension': ext,
        'extracted_text': None,
        'error': None
    }
    
    try:
        if ext == '.pdf':
            result['extracted_text'] = extract_text_from_pdf(file_path)
        elif ext == '.docx':
            result['extracted_text'] = extract_text_from_docx(file_path)
        elif ext == '.xlsx':
            result['extracted_text'] = extract_text_from_xlsx(file_path)
        elif ext == '.pptx':
            result['extracted_text'] = extract_text_from_pptx(file_path)
        elif ext == '.txt':
            result['extracted_text'] = extract_text_from_txt(file_path)
        else:
            result['error'] = f"Unsupported file type: {ext}"
            return result
        
        # Truncate if too long (optional, to avoid huge outputs)
        if result['extracted_text'] and len(result['extracted_text']) > 100000:
            result['extracted_text'] = result['extracted_text'][:100000] + "\n\n[TRUNCATED: Output exceeds 100k characters]"
            result['truncated'] = True
        
    except Exception as e:
        result['error'] = str(e)
        logger.error(f"Error processing {file_path}: {e}")
    
    return result


def main():
    parser = argparse.ArgumentParser(description='Universal Analyzer - Extract text from documents')
    parser.add_argument('file', help='Path to the file to analyze')
    parser.add_argument('--output', '-o', help='Output file to save extracted text (optional)')
    parser.add_argument('--quiet', '-q', action='store_true', help='Suppress verbose output')
    parser.add_argument('--metadata', '-m', action='store_true', help='Show file metadata only')
    
    args = parser.parse_args()
    
    if not args.quiet:
        logger.info(f"Analyzing: {args.file}")
    
    try:
        result = analyze_file(args.file)
        
        if result['error']:
            print(f"ERROR: {result['error']}", file=sys.stderr)
            sys.exit(1)
        
        if args.metadata:
            print(f"File: {result['file_name']}")
            print(f"Size: {result['file_size']} bytes")
            print(f"Type: {result['extension']}")
            if result.get('truncated'):
                print("Warning: Output was truncated")
        else:
            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    f.write(result['extracted_text'])
                if not args.quiet:
                    logger.info(f"Text saved to: {args.output}")
            else:
                print(result['extracted_text'])
    
    except Exception as e:
        logger.error(f"Failed to analyze file: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
